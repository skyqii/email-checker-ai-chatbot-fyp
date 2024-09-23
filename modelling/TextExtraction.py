from pdfminer.high_level import extract_text
import os
from docx import Document
from pptx import Presentation
import openpyxl

# Function to extract text from PDF files using pdfminer.six
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        text += extract_text(pdf)
    return text

# Function to extract text from DOCX files including text from images using OCR
def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        
        # Extracting text from paragraphs
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        # Extracting text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"  # Newline after each row
    
    return text

def get_pptx_text(pptx_docs):
    text = ""
    
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        
        # Loop through each slide in the presentation
        for slide in presentation.slides:

            # Extract text from shapes (text boxes, titles, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            # Extract text from speaker notes (if any)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide:
                    notes_text = notes_slide.notes_text_frame.text
                    text += f"{notes_text}\n"
                
    return text

# Function to extract text from XLSX files
def get_xlsx_text(xlsx_docs):
    text = ""
    for xlsx in xlsx_docs:
        wb = openpyxl.load_workbook(xlsx)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text += str(cell.value) + " "
                text += "\n"
    return text

# Directory paths
pdf_dir = r'modelling\data\pdf_files'
docx_dir = r'modelling\data\docx_files'
pptx_dir = r'modelling\data\pptx_files'
xlsx_dir = r'modelling\data\xlsx_files'

# Fetch file paths
pdf_files = [os.path.join(pdf_dir, f) for f in os.listdir(pdf_dir) if f.endswith('.pdf')]
docx_files = [os.path.join(docx_dir, f) for f in os.listdir(docx_dir) if f.endswith('.docx')]
pptx_files = [os.path.join(pptx_dir, f) for f in os.listdir(pptx_dir) if f.endswith('.pptx')]
xlsx_files = [os.path.join(xlsx_dir, f) for f in os.listdir(xlsx_dir) if f.endswith('.xlsx')]

# Extract text for each document type
pdf_text = get_pdf_text(pdf_files)
docx_text = get_docx_text(docx_files)
pptx_text = get_pptx_text(pptx_files)
xlsx_text = get_xlsx_text(xlsx_files)

# Define the output file for all extracted text
output_file = os.path.join('modelling\data', 'extracted_text.txt')

# Function to store all extracted text into one file
def save_all_text_to_file(texts, file_name):
    with open(file_name, 'w', encoding='utf-8') as f:
        for text_type, text in texts.items():
            # f.write(f"--- {text_type.upper()} TEXT ---\n")
            f.write(text + "\n")

# Combine all extracted texts into a dictionary
all_texts = {
    "pdf": pdf_text,
    "docx": docx_text,
    "pptx": pptx_text,
    "xlsx": xlsx_text
}

# Save all text into one file
save_all_text_to_file(all_texts, output_file)
